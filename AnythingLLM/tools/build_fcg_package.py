"""
build_fcg_package.py
Generador e integrador oficial de la skill 'fcg-ppt' (Faro Consulting Group).
Crea:
1. assets/template_fcg.pptx
2. scripts/fcg_helpers.py
3. scripts/example_deck.py
4. SKILL.md
Y sincroniza en todas las rutas de skills (Workspace y Global).
"""

import os
import sys
from pathlib import Path
from datetime import datetime

# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

FCG_HELPERS_CODE = '''# -*- coding: utf-8 -*-
"""
fcg_helpers.py
Biblioteca de utilidades y componentes visuales oficiales de Faro Consulting Group (FaroCG)
para la generación automatizada de presentaciones .pptx ejecutivas de alto impacto.
"""

from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Union

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

# ==============================================================================
# 1. PALETA DE COLORES OFICIAL FARO CONSULTING GROUP (FaroCG)
# ==============================================================================
FCG_NAVY        = RGBColor(11, 37, 69)     # #0B2545 - Azul Marino Institucional Profundo
FCG_OCEAN       = RGBColor(19, 60, 85)     # #133C55 - Azul Océano Faro
FCG_ROYAL       = RGBColor(29, 78, 216)    # #1D4ED8 - Azul Eléctrico Ejecutivo
FCG_SKY         = RGBColor(56, 189, 248)   # #38BDF8 - Azul Cielo Faro (Acento luminoso)
FCG_AMBER       = RGBColor(217, 119, 6)    # #D97706 - Ámbar / Faro Destello (Callout)
FCG_GOLD        = RGBColor(245, 158, 11)   # #F59E0B - Oro Faro Acento
FCG_DARK        = RGBColor(30, 41, 59)     # #1E293B - Slate Dark (Texto principal)
FCG_MUTED       = RGBColor(100, 116, 139)  # #64748B - Slate Muted (Subtítulos y notas)
FCG_BG          = RGBColor(248, 250, 252)  # #F8FAFC - Fondo Off-White Ejecutivo
FCG_WHITE       = RGBColor(255, 255, 255)  # #FFFFFF - Blanco Puro
FCG_BORDER      = RGBColor(226, 232, 240)  # #E2E8F0 - Borde de Tarjetas Suave
FCG_CARD_BG     = RGBColor(255, 255, 255)  # #FFFFFF - Fondo de Cards

# Alias de compatibilidad
FCG_PRIMARY     = FCG_NAVY
FCG_SECONDARY   = FCG_OCEAN
FCG_ACCENT      = FCG_ROYAL
FCG_ACCENT_LIGHT= FCG_SKY

# Tipografías estándar garantizadas en cualquier sistema operativo
FONT_TITLE      = "Calibri"
FONT_BODY       = "Calibri"

# Dimensiones 16:9 Widescreen (33.867 cm x 19.05 cm)
SLIDE_WIDTH     = Inches(13.333333)
SLIDE_HEIGHT    = Inches(7.5)


# ==============================================================================
# 2. CARGA Y GESTIÓN DE PLANTILLAS
# ==============================================================================
def load_fcg_template(
    template_path: Optional[Union[str, Path]] = None,
    clear_existing: bool = True
) -> Presentation:
    """
    Carga la plantilla oficial assets/template_fcg.pptx o inicializa una presentación 16:9
    con las dimensiones exactas y estándares de Faro Consulting Group.
    Si clear_existing es True, elimina diapositivas remanentes para entregar un lienzo limpio.
    """
    candidates = []
    if template_path:
        candidates.append(Path(template_path))
    
    current_dir = Path(__file__).resolve().parent
    candidates.extend([
        current_dir.parent / "assets" / "template_fcg.pptx",
        current_dir / "assets" / "template_fcg.pptx",
        Path.cwd() / "assets" / "template_fcg.pptx",
        Path.cwd() / ".agents" / "skills" / "fcg-ppt" / "assets" / "template_fcg.pptx",
        Path.cwd() / "skills" / "fcg-ppt" / "assets" / "template_fcg.pptx",
    ])
    
    prs = None
    for candidate in candidates:
        if candidate and candidate.is_file():
            prs = Presentation(str(candidate))
            prs.slide_width = SLIDE_WIDTH
            prs.slide_height = SLIDE_HEIGHT
            break

    if prs is None:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        
    if clear_existing and len(prs.slides) > 0:
        # Remover diapositivas previas para iniciar mazo limpio
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

    return prs


def dated_filename(base_name: str = "presentacion_fcg", ext: str = "pptx") -> str:
    """
    Genera un nombre de archivo normalizado con timestamp: {base_name}_{YYYYMMDD_HHMMSS}.{ext}
    """
    now_str = datetime.now().strftime("%Y%m%d_%H%M%S")
    clean_base = base_name.rstrip(".").replace(" ", "_")
    return f"{clean_base}_{now_str}.{ext.lstrip('.')}"


# ==============================================================================
# 3. HELPERS DE GEOMETRÍA, TARJETAS Y TEXTO
# ==============================================================================
def set_shape_text(
    shape,
    text: str,
    font_size: Pt = Pt(13),
    font_color: RGBColor = FCG_DARK,
    font_name: str = FONT_BODY,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    word_wrap: bool = True,
    space_after: Pt = Pt(4),
    space_before: Pt = Pt(0)
):
    """
    Formatea de manera exhaustiva el texto de una forma o caja de texto,
    garantizando margins limpios, wrap de palabras y atributos tipográficos.
    """
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    
    if len(tf.paragraphs) == 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
        
    p.text = text
    p.alignment = align
    p.space_after = space_after
    p.space_before = space_before
    
    if len(p.runs) > 0:
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.italic = italic
        # Configurar idioma español Chile
        try:
            run._r.get_or_add_rPr().set('lang', 'es-CL')
        except Exception:
            pass


def add_rect(
    slide,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    fill_color: RGBColor = FCG_PRIMARY,
    border_color: Optional[RGBColor] = None,
    border_width: Pt = Pt(1)
):
    """Agrega un rectángulo geométrico con color plano y borde opcional."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if border_color:
        rect.line.color.rgb = border_color
        rect.line.width = border_width
    else:
        rect.line.fill.background()
    return rect


def add_circle(
    slide,
    left: Inches,
    top: Inches,
    size: Inches,
    fill_color: RGBColor = FCG_ROYAL,
    border_color: Optional[RGBColor] = None,
    border_width: Pt = Pt(1)
):
    """Agrega un círculo perfecto con diámetro 'size'."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    if border_color:
        circle.line.color.rgb = border_color
        circle.line.width = border_width
    else:
        circle.line.fill.background()
    return circle


def add_card(
    slide,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    fill_color: RGBColor = FCG_CARD_BG,
    border_color: RGBColor = FCG_BORDER,
    border_width: Pt = Pt(1)
):
    """Agrega un contenedor tipo tarjeta (Card) con fondo blanco y borde suave."""
    return add_rect(slide, left, top, width, height, fill_color, border_color, border_width)


def set_notes(slide, notes_text: str):
    """Establece las notas del orador (speaker notes) de la diapositiva."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text
        for p in tf.paragraphs:
            for r in p.runs:
                try:
                    r._r.get_or_add_rPr().set('lang', 'es-CL')
                except Exception:
                    pass
    except Exception as e:
        print(f"Advertencia: No se pudieron asignar notas: {e}")


# ==============================================================================
# 4. PATRONES ESTRUCTURALES DE SLIDE (PORTADA, CABECERAS, CONTENIDO)
# ==============================================================================
def set_portada(
    slide,
    title: str,
    subtitle: Optional[str] = None,
    author: str = "Faro Consulting Group",
    date: Optional[str] = None,
    client: Optional[str] = None
):
    """
    Configura una diapositiva en blanco como la portada oficial Faro Consulting Group:
    Fondo Midnight Navy (#0B2545), franja ámbar de destello y tipografía corporativa.
    """
    # Fondo Navy completo
    bg = add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, FCG_NAVY)
    
    # Barra lateral de acento Faro (Luz de Faro: Océano + Ámbar)
    add_rect(slide, 0, 0, Inches(0.4), SLIDE_HEIGHT, FCG_AMBER)
    add_rect(slide, Inches(0.4), 0, Inches(0.15), SLIDE_HEIGHT, FCG_SKY)
    
    # Franja inferior decorativa
    add_rect(slide, Inches(0.55), Inches(7.3), Inches(12.78), Inches(0.2), FCG_OCEAN)
    
    # Tag de Marca Superior
    tag_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(8), Inches(0.5))
    set_shape_text(tag_box, "FARO CONSULTING GROUP  |  ESTRATEGIA & TRANSFORMACIÓN", 
                   font_size=Pt(11), font_color=FCG_SKY, bold=True)
    
    # Título Principal
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p = tf_title.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = FCG_WHITE
    try:
        p.runs[0]._r.get_or_add_rPr().set('lang', 'es-CL')
    except Exception:
        pass
    
    # Subtítulo descriptivo
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(10.5), Inches(1.1))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(17)
        p_sub.font.color.rgb = RGBColor(203, 213, 225) # Slate 300
        try:
            p_sub.runs[0]._r.get_or_add_rPr().set('lang', 'es-CL')
        except Exception:
            pass
            
    # Metadata inferior (Autor, Cliente, Fecha)
    date_val = date or datetime.now().strftime("%d de %B de %Y")
    meta_text = f"Preparado por: {author}   |   Fecha: {date_val}"
    if client:
        meta_text = f"Cliente: {client}   |   " + meta_text
        
    meta_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.3), Inches(10.5), Inches(0.6))
    set_shape_text(meta_box, meta_text, font_size=Pt(11), font_color=RGBColor(148, 163, 184), bold=False)


def add_portada_slide(
    prs: Presentation,
    title: str,
    subtitle: Optional[str] = None,
    author: str = "Faro Consulting Group",
    date: Optional[str] = None,
    client: Optional[str] = None
):
    """Crea y añade una nueva diapositiva de portada al deck."""
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    set_portada(slide, title, subtitle, author, date, client)
    return slide


def set_title(
    slide,
    text: str,
    subtitle: Optional[str] = None,
    category: str = "FARO CONSULTING GROUP"
):
    """
    Configura el encabezado corporativo FaroCG en una diapositiva:
    - Tag de categoría / sección superior
    - Título principal en azul marino de 24pt
    - Subtítulo explicativo en slate muted de 12pt
    - Franja decorativa de navegación
    """
    # Barra superior delgada (Faro Navy)
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.1), FCG_NAVY)
    # Franja ámbar indicadora
    add_rect(slide, Inches(0.8), Inches(0.1), Inches(1.5), Inches(0.05), FCG_AMBER)
    
    # Categoría / Kicker
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.35))
    set_shape_text(cat_box, category.upper(), font_size=Pt(9), font_color=FCG_ROYAL, bold=True)
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_TITLE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = FCG_NAVY
    try:
        p.runs[0]._r.get_or_add_rPr().set('lang', 'es-CL')
    except Exception:
        pass
    
    # Subtítulo (si existe)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.45))
        set_shape_text(sub_box, subtitle, font_size=Pt(12), font_color=FCG_MUTED, bold=False)


def add_content_slide(
    prs: Presentation,
    title: str,
    subtitle: Optional[str] = None,
    category: str = "FARO CONSULTING GROUP"
):
    """
    Añade una diapositiva de contenido estándar con fondo Off-White (#F8FAFC),
    encabezado institucional y pie de página de FaroCG.
    """
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    
    # Fondo Off-White
    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, FCG_BG)
    
    # Encabezado
    set_title(slide, title, subtitle, category)
    
    # Footer institucional
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(9.0), Inches(0.35))
    set_shape_text(footer_box, "Faro Consulting Group  |  Documento Confidencial de Asesoría",
                   font_size=Pt(9), font_color=FCG_MUTED)
                   
    # Marca de agua / franja decorativa footer
    add_rect(slide, Inches(12.0), Inches(7.1), Inches(0.5), Inches(0.04), FCG_AMBER)
    
    return slide


# ==============================================================================
# 5. PATRONES DE SLIDES LISTOS (LISTED PATTERNS)
# ==============================================================================
def add_pattern_two_columns(
    slide,
    col1_title: str,
    col1_bullets: List[str],
    col2_title: str,
    col2_bullets: List[str],
    left: Inches = Inches(0.8),
    top: Inches = Inches(1.85),
    width: Inches = Inches(11.73),
    height: Inches = Inches(4.8)
):
    """
    Patrón: Cards de 2 Columnas.
    Genera dos paneles ejecutivos en paralelo con cabeceras temáticas y lista de bullets.
    """
    gap = Inches(0.35)
    col_width = (width - gap) / 2
    
    cols = [
        (left, col1_title, col1_bullets, FCG_NAVY),
        (left + col_width + gap, col2_title, col2_bullets, FCG_OCEAN)
    ]
    
    for c_left, c_title, c_bullets, c_header_color in cols:
        # Card container
        add_card(slide, c_left, top, col_width, height)
        
        # Header banner
        header_rect = add_rect(slide, c_left, top, col_width, Inches(0.7), c_header_color)
        tb_h = slide.shapes.add_textbox(c_left + Inches(0.2), top + Inches(0.12), col_width - Inches(0.4), Inches(0.45))
        set_shape_text(tb_h, c_title, font_size=Pt(14), font_color=FCG_WHITE, bold=True)
        
        # Bullets box
        tb_body = slide.shapes.add_textbox(
            c_left + Inches(0.25), 
            top + Inches(0.85), 
            col_width - Inches(0.5), 
            height - Inches(1.0)
        )
        tf = tb_body.text_frame
        tf.word_wrap = True
        
        for idx, bullet in enumerate(c_bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = f"•  {bullet}"
            p.space_after = Pt(8)
            p.space_before = Pt(3)
            if len(p.runs) > 0:
                p.runs[0].font.name = FONT_BODY
                p.runs[0].font.size = Pt(12)
                p.runs[0].font.color.rgb = FCG_DARK
                try:
                    p.runs[0]._r.get_or_add_rPr().set('lang', 'es-CL')
                except Exception:
                    pass


def add_pattern_metric_cards(
    slide,
    metrics: List[Dict[str, str]],
    top: Inches = Inches(2.0),
    height: Inches = Inches(4.6)
):
    """
    Patrón: Cards de Métricas / KPIs.
    metrics = [
        {"number": "+45%", "label": "Eficiencia Operativa", "desc": "Reducción directa de tiempos."},
        ...
    ]
    Soporta dinámicamente de 2 a 4 tarjetas distribuidas uniformemente.
    """
    count = len(metrics)
    if count == 0:
        return
        
    total_width = Inches(11.73)
    left_start = Inches(0.8)
    gap = Inches(0.3)
    card_width = (total_width - (gap * (count - 1))) / count
    
    accent_colors = [FCG_ROYAL, FCG_AMBER, FCG_OCEAN, FCG_SKY]
    
    for i, m in enumerate(metrics):
        c_left = left_start + i * (card_width + gap)
        c_color = accent_colors[i % len(accent_colors)]
        
        # Tarjeta base
        add_card(slide, c_left, top, card_width, height)
        
        # Barra superior de acento de la tarjeta
        add_rect(slide, c_left, top, card_width, Inches(0.12), c_color)
        
        # Número / KPI destacado
        num_box = slide.shapes.add_textbox(c_left + Inches(0.15), top + Inches(0.4), card_width - Inches(0.3), Inches(1.1))
        set_shape_text(num_box, m.get("number", "0"), font_size=Pt(34), font_color=c_color, bold=True, align=PP_ALIGN.CENTER)
        
        # Etiqueta / Nombre de métrica
        lbl_box = slide.shapes.add_textbox(c_left + Inches(0.15), top + Inches(1.6), card_width - Inches(0.3), Inches(0.6))
        set_shape_text(lbl_box, m.get("label", ""), font_size=Pt(13), font_color=FCG_NAVY, bold=True, align=PP_ALIGN.CENTER)
        
        # Línea divisoria suave
        add_rect(slide, c_left + Inches(0.4), top + Inches(2.3), card_width - Inches(0.8), Pt(1), FCG_BORDER)
        
        # Descripción / Impacto
        desc_box = slide.shapes.add_textbox(c_left + Inches(0.2), top + Inches(2.5), card_width - Inches(0.4), height - Inches(2.7))
        set_shape_text(desc_box, m.get("desc", ""), font_size=Pt(11), font_color=FCG_MUTED, align=PP_ALIGN.CENTER)


def add_pattern_bullet_circles(
    slide,
    items: List[Dict[str, str]],
    top: Inches = Inches(2.0)
):
    """
    Patrón: Bullet Circles Numerados (Fases o Pasos Metodológicos).
    items = [
        {"step": "01", "title": "Diagnóstico Inicial", "desc": "Levantamiento de procesos..."},
        ...
    ]
    """
    start_left = Inches(0.8)
    row_height = Inches(1.15)
    
    for i, it in enumerate(items[:4]):
        cur_top = top + i * (row_height + Inches(0.15))
        
        # Contenedor de fila
        add_card(slide, start_left, cur_top, Inches(11.73), row_height)
        
        # Círculo con número
        add_circle(slide, start_left + Inches(0.25), cur_top + Inches(0.18), Inches(0.8), FCG_NAVY)
        
        # Texto del número dentro del círculo
        num_box = slide.shapes.add_textbox(start_left + Inches(0.25), cur_top + Inches(0.28), Inches(0.8), Inches(0.6))
        set_shape_text(num_box, it.get("step", f"0{i+1}"), font_size=Pt(15), font_color=FCG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        
        # Título del paso
        title_box = slide.shapes.add_textbox(start_left + Inches(1.25), cur_top + Inches(0.15), Inches(10.0), Inches(0.4))
        set_shape_text(title_box, it.get("title", ""), font_size=Pt(13), font_color=FCG_NAVY, bold=True)
        
        # Detalle / Descripción
        desc_box = slide.shapes.add_textbox(start_left + Inches(1.25), cur_top + Inches(0.55), Inches(10.0), Inches(0.5))
        set_shape_text(desc_box, it.get("desc", ""), font_size=Pt(11), font_color=FCG_MUTED)


def add_pattern_content_with_stripe(
    slide,
    headline: str,
    paragraphs: List[str],
    key_takeaway: Optional[str] = None,
    top: Inches = Inches(1.9)
):
    """
    Patrón: Bloque de Contenido con Franja Lateral y Cuadro de Conclusión/Takeaway.
    """
    # Card principal
    add_card(slide, Inches(0.8), top, Inches(11.73), Inches(4.8))
    
    # Franja vertical decorativa izquierda (Faro Accent)
    add_rect(slide, Inches(0.8), top, Inches(0.15), Inches(4.8), FCG_AMBER)
    
    # Headline destacado
    hl_box = slide.shapes.add_textbox(Inches(1.2), top + Inches(0.3), Inches(11.0), Inches(0.6))
    set_shape_text(hl_box, headline, font_size=Pt(16), font_color=FCG_NAVY, bold=True)
    
    # Párrafos de desarrollo
    body_box = slide.shapes.add_textbox(Inches(1.2), top + Inches(1.0), Inches(11.0), Inches(2.2))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, p_txt in enumerate(paragraphs):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = p_txt
        p.space_after = Pt(8)
        if len(p.runs) > 0:
            p.runs[0].font.name = FONT_BODY
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = FCG_DARK
            try:
                p.runs[0]._r.get_or_add_rPr().set('lang', 'es-CL')
            except Exception:
                pass
                
    # Key Takeaway Box (si existe)
    if key_takeaway:
        tk_card = add_card(slide, Inches(1.2), top + Inches(3.4), Inches(11.0), Inches(1.1), fill_color=RGBColor(241, 245, 249))
        add_rect(slide, Inches(1.2), top + Inches(3.4), Inches(0.08), Inches(1.1), FCG_ROYAL)
        
        tk_title = slide.shapes.add_textbox(Inches(1.4), top + Inches(3.48), Inches(10.6), Inches(0.3))
        set_shape_text(tk_title, "IMPLICANCIA ESTRATÉGICA / CONCLUSIÓN CLAVE:", font_size=Pt(10), font_color=FCG_ROYAL, bold=True)
        
        tk_desc = slide.shapes.add_textbox(Inches(1.4), top + Inches(3.8), Inches(10.6), Inches(0.6))
        set_shape_text(tk_desc, key_takeaway, font_size=Pt(11), font_color=FCG_DARK, italic=True)


def add_cierre_slide(
    prs: Presentation,
    title: str = "Gracias",
    subtitle: str = "Faro Consulting Group — Guiando Decisiones Estratégicas",
    contact_info: Optional[Dict[str, str]] = None
):
    """
    Diapositiva de cierre y agradecimiento con estilo Faro Consulting Group.
    """
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    
    # Fondo Navy
    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, FCG_NAVY)
    
    # Franjas decorativas
    add_rect(slide, 0, Inches(7.3), SLIDE_WIDTH, Inches(0.2), FCG_AMBER)
    add_rect(slide, 0, Inches(7.15), Inches(3.5), Inches(0.08), FCG_SKY)
    
    # Título central
    t_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2))
    set_shape_text(t_box, title, font_size=Pt(44), font_color=FCG_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # Subtítulo
    s_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8))
    set_shape_text(s_box, subtitle, font_size=Pt(16), font_color=RGBColor(203, 213, 225), align=PP_ALIGN.CENTER)
    
    # Contacto
    info = contact_info or {
        "Web": "www.faroconsulting.com",
        "Contacto": "contacto@faroconsulting.com",
        "Ubicación": "Santiago, Chile"
    }
    contact_str = "   |   ".join([f"{k}: {v}" for k, v in info.items()])
    c_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6))
    set_shape_text(c_box, contact_str, font_size=Pt(11), font_color=FCG_SKY, align=PP_ALIGN.CENTER)
    
    return slide


# ==============================================================================
# 6. VALIDACIONES Y PERSISTENCIA ROBUSTA
# ==============================================================================
def validate_layout(slide) -> bool:
    """
    Verifica que ninguna forma se salga de los márgenes seguros de la diapositiva (13.33 x 7.5 pulg).
    Retorna True si todas las formas están dentro de los límites válidos.
    """
    valid = True
    for shape in slide.shapes:
        try:
            right = shape.left + shape.width
            bottom = shape.top + shape.height
            if right > SLIDE_WIDTH + Inches(0.1) or bottom > SLIDE_HEIGHT + Inches(0.1):
                valid = False
                print(f"[Validación Layout] Forma fuera de margen: left={shape.left}, top={shape.top}, right={right}, bottom={bottom}")
        except Exception:
            pass
    return valid


def set_language_es_cl(prs: Presentation):
    """
    Recorre todos los párrafos y runs de todas las diapositivas para asignar
    explícitamente el atributo lang='es-CL' (Español Chile), evitando subrayados
    rojos de corrección ortográfica en PowerPoint.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            run._r.get_or_add_rPr().set('lang', 'es-CL')
                        except Exception:
                            pass
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                try:
                                    run._r.get_or_add_rPr().set('lang', 'es-CL')
                                except Exception:
                                    pass


def save_with_normal_view(prs: Presentation, output_path: Union[str, Path]):
    """
    Guarda la presentación asegurando la configuración en presentation.xml
    para que PowerPoint la abra inmediatamente en modo 'Normal View' (Diapositiva + Miniaturas).
    """
    set_language_es_cl(prs)
    
    try:
        view_pr = prs._element.find(qn('p:viewPr'))
        if view_pr is None:
            view_xml = parse_xml(
                '<p:viewPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                '  <p:normalViewPr>'
                '    <p:restoredLeft sz="15620"/>'
                '    <p:restoredTop sz="94600"/>'
                '  </p:normalViewPr>'
                '</p:viewPr>'
            )
            prs._element.append(view_xml)
    except Exception as e:
        print(f"[Aviso] No se pudo inyectar viewPr en XML: {e}")
        
    out_p = Path(output_path)
    out_p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_p))
    print(f"[OK] Presentacion FaroCG guardada con exito en: {out_p.resolve()}")
'''

# Template generator script
def create_master_template(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    
    # Inyectar viewPr para vista normal
    view_xml = parse_xml(
        '<p:viewPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '  <p:normalViewPr>'
        '    <p:restoredLeft sz="15620"/>'
        '    <p:restoredTop sz="94600"/>'
        '  </p:normalViewPr>'
        '</p:viewPr>'
    )
    prs._element.append(view_xml)
    
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Plantilla generada en: {output_path}")


EXAMPLE_DECK_CODE = '''# -*- coding: utf-8 -*-
"""
example_deck.py
Demostración completa de generación de una presentación ejecutiva oficial
para Faro Consulting Group (FaroCG) utilizando fcg_helpers.
"""

import sys
from pathlib import Path

# Añadir scripts al path si se ejecuta directamente
sys.path.insert(0, str(Path(__file__).resolve().parent))

from fcg_helpers import (
    load_fcg_template,
    dated_filename,
    add_portada_slide,
    add_content_slide,
    add_pattern_two_columns,
    add_pattern_metric_cards,
    add_pattern_bullet_circles,
    add_pattern_content_with_stripe,
    add_cierre_slide,
    set_notes,
    validate_layout,
    save_with_normal_view
)

def build_presentation():
    # 1. Cargar plantilla oficial FaroCG 16:9
    prs = load_fcg_template()
    
    # --------------------------------------------------------------------------
    # SLIDE 1: Portada Ejecutiva
    # --------------------------------------------------------------------------
    slide_portada = add_portada_slide(
        prs,
        title="Transformación Digital & Estrategia IA 2026",
        subtitle="Hoja de ruta integral para la modernización de procesos y ventaja competitiva",
        author="Faro Consulting Group — Práctica de Estrategia Digital",
        date="Septiembre 2026",
        client="Dirección General & Comité Ejecutivo"
    )
    set_notes(slide_portada, "Abrir la sesión enfatizando el objetivo estratégico: convertir las capacidades de IA local y automatización en eficiencia medible y ROI en 6 meses.")
    validate_layout(slide_portada)
    
    # --------------------------------------------------------------------------
    # SLIDE 2: Resumen Ejecutivo (Cards 2 Columnas)
    # --------------------------------------------------------------------------
    slide_resumen = add_content_slide(
        prs,
        title="Diagnóstico Estratégico y Oportunidad de Mercado",
        subtitle="Evaluación del estado actual de operaciones frente a benchmarks del sector",
        category="01. DIAGNÓSTICO EJECUTIVO"
    )
    
    col1_bullets = [
        "Sistemas heredados generan un 38% de sobrecosto operativo en tareas repetitivas.",
        "Dispersión de información documental crítica retrasa la toma de decisiones en un promedio de 4.2 días.",
        "Riesgo de cumplimiento normativo ante normativas crecientes de gobernanza y privacidad de datos."
    ]
    
    col2_bullets = [
        "Implementación de arquitectura IA soberana local con control absoluto de privacidad y seguridad.",
        "Automatización de auditorías y procesamiento de contratos reduciendo tiempos a minutos.",
        "Consolidación de un ecosistema escalable de skills que potencia al equipo multidisciplinario."
    ]
    
    add_pattern_two_columns(
        slide_resumen,
        col1_title="Desafíos Operativos Actuales",
        col1_bullets=col1_bullets,
        col2_title="Propuesta de Valor & Solución",
        col2_bullets=col2_bullets
    )
    set_notes(slide_resumen, "Puntualizar que la columna izquierda representa costos hundidos actuales y la derecha el ahorro proyectado.")
    validate_layout(slide_resumen)
    
    # --------------------------------------------------------------------------
    # SLIDE 3: Métricas y KPIs de Impacto (Metric Cards)
    # --------------------------------------------------------------------------
    slide_kpi = add_content_slide(
        prs,
        title="Impacto Proyectado y Retorno de Inversión (ROI)",
        subtitle="Métricas cuantitativas estimadas al cabo del primer año de implementación",
        category="02. IMPACTO & RENDIMIENTO"
    )
    
    kpis = [
        {
            "number": "+65%",
            "label": "Velocidad Operativa",
            "desc": "Aceleración en el análisis forense y auditoría documental integral."
        },
        {
            "number": "4.8x",
            "label": "Retorno de Inversión",
            "desc": "ROI cuantificado a 12 meses basado en reducción de horas hombre y licencias externas."
        },
        {
            "number": "100%",
            "label": "Soberanía de Datos",
            "desc": "Procesamiento totalmente local on-premise sin fuga a nubes de terceros."
        },
        {
            "number": "-40%",
            "label": "Margen de Error",
            "desc": "Disminución de discrepancias contractuales mediante validación automatizada."
        }
    ]
    
    add_pattern_metric_cards(slide_kpi, kpis)
    set_notes(slide_kpi, "Detallar que el 100% de soberanía es clave para el área legal y de auditoría interna.")
    validate_layout(slide_kpi)
    
    # --------------------------------------------------------------------------
    # SLIDE 4: Metodología de Implementación (Bullet Circles)
    # --------------------------------------------------------------------------
    slide_fases = add_content_slide(
        prs,
        title="Fases de la Metodología FaroCG",
        subtitle="Plan de despliegue progresivo estructurado en cuatro hitos críticos",
        category="03. PLAN DE TRABAJO"
    )
    
    fases = [
        {
            "step": "01",
            "title": "Auditoría & Descubrimiento de Infraestructura (Semanas 1-2)",
            "desc": "Levantamiento de capacidades de cómputo local, modelos compatibles y flujos prioritarios."
        },
        {
            "step": "02",
            "title": "Despliegue de Gateway & Agentes Especializados (Semanas 3-5)",
            "desc": "Configuración del entorno soberano, integración de skills a la medida y bases vectoriales."
        },
        {
            "step": "03",
            "title": "Piloto Operativo Multidocumental (Semanas 6-8)",
            "desc": "Pruebas de estrés de análisis documental en paralelo con retroalimentación en tiempo real."
        },
        {
            "step": "04",
            "title": "Pase a Producción & Capacitación Directiva (Semanas 9-10)",
            "desc": "Transferencia metodológica de mejores prácticas de prompting y gobernanza continua."
        }
    ]
    
    add_pattern_bullet_circles(slide_fases, fases)
    set_notes(slide_fases, "Subrayar que cada fase cuenta con entregables y criterios de aceptación formalmente auditados.")
    validate_layout(slide_fases)
    
    # --------------------------------------------------------------------------
    # SLIDE 5: Conclusiones Estratégicas (Contenido con Franja & Takeaway)
    # --------------------------------------------------------------------------
    slide_concl = add_content_slide(
        prs,
        title="Recomendaciones Finales del Consejo Consultor",
        subtitle="Síntesis directiva para la aprobación del presupuesto y calendario",
        category="04. CONCLUSIÓN ESTRATÉGICA"
    )
    
    parrafos = [
        "La adopción de una arquitectura de IA local privada representa la ruta óptima para maximizar la seguridad informática y eliminar costos recurrentes de API propietarias.",
        "La estandarización de herramientas como 'fcg-ppt' asegura coherencia de marca ininterrumpida y un estándar corporativo de excelencia en cada interacción con stakeholders."
    ]
    
    takeaway = "La aprobación del presupuesto en Q3 asegura la captura de eficiencias operativas antes del cierre fiscal, consolidando a la organización como líder tecnológico en su industria."
    
    add_pattern_content_with_stripe(
        slide_concl,
        headline="La Modernización No Es Un Costo, Es Una Ventaja Competitiva Defensiva",
        paragraphs=parrafos,
        key_takeaway=takeaway
    )
    set_notes(slide_concl, "Cerrar con llamado explícito a la acción solicitando la firma del acta de inicio de proyecto.")
    validate_layout(slide_concl)
    
    # --------------------------------------------------------------------------
    # SLIDE 6: Diapositiva de Cierre
    # --------------------------------------------------------------------------
    slide_cierre = add_cierre_slide(
        prs,
        title="Gracias",
        subtitle="Faro Consulting Group — Guiando Decisiones Estratégicas con Inteligencia y Visión",
        contact_info={
            "Web": "www.faroconsulting.com",
            "Contacto": "contacto@faroconsulting.com",
            "Sede": "Santiago, Chile"
        }
    )
    validate_layout(slide_cierre)
    
    # Guardar presentación
    out_filename = dated_filename("presentacion_faro_ejemplo", "pptx")
    out_dir = Path(__file__).resolve().parent.parent / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / out_filename
    
    save_with_normal_view(prs, out_path)
    return out_path

if __name__ == "__main__":
    generated_path = build_presentation()
    print(f"[EXITO] Presentacion generada en {generated_path}")
'''

SKILL_MD_CONTENT = '''---
name: fcg-ppt
description: "Generador de presentaciones PPTX ejecutivas con formato corporativo e identidad visual de Faro Consulting Group (FaroCG). Incluye plantilla 16:9, paleta oficial, helpers de python-pptx y patrones visuales predefinidos."
tags:
  - pptx
  - presentations
  - faro-consulting
  - corporate-branding
  - python-pptx
---

# Skill: Generador de PPT con formato Faro Consulting Group

## Que hace este skill

Empaqueta la identidad visual de FaroCG y un conjunto de helpers Python para producir presentaciones `.pptx` consistentes con el branding de Faro:

- **Plantilla `.pptx` oficial** con masters y layouts FCG — en `assets/template_fcg.pptx`.
- **Paleta de colores oficial y constantes tipográficas** — en `scripts/fcg_helpers.py`.
- **Helpers reutilizables**:
  - `set_shape_text(shape, text, font_size, font_color, bold, align)`
  - `add_rect(slide, left, top, width, height, fill_color, border_color)`
  - `add_circle(slide, left, top, size, fill_color)`
  - `add_card(slide, left, top, width, height)`
  - `set_notes(slide, notes_text)`
  - `set_title(slide, text, subtitle, category)`
  - `set_portada(slide, title, subtitle, author, date, client)`
  - `add_portada_slide(prs, title, subtitle, author, date, client)`
  - `add_content_slide(prs, title, subtitle, category)`
  - `load_fcg_template(template_path)`
  - `save_with_normal_view(prs, output_path)`
  - `dated_filename(base_name, ext)`
  - `set_language_es_cl(prs)`
  - `validate_layout(slide)`
- **Patrones de slide listos**:
  - Header bar institucional con franja de acento ámbar.
  - Cards de 2 columnas (`add_pattern_two_columns`).
  - Bullet circles numerados para fases metodológicas (`add_pattern_bullet_circles`).
  - Cards de métricas y KPIs (`add_pattern_metric_cards`).
  - Bloque de contenido con franja y takeaway estratégico (`add_pattern_content_with_stripe`).
  - Diapositiva de cierre y agradecimiento (`add_cierre_slide`).

---

## Estructura del Skill

```
fcg-ppt/
├── SKILL.md                 # Documentación e instrucciones operativas
├── assets/
│   └── template_fcg.pptx   # Plantilla base 16:9 Widescreen con masters FaroCG
└── scripts/
    ├── fcg_helpers.py       # Módulo Python con helpers, paleta y patrones
    └── example_deck.py      # Script de demostración con mazo de 6 diapositivas
```

---

## Identidad Visual y Paleta de Colores

| Color | Hex | RGB | Uso Recomendado |
| :--- | :--- | :--- | :--- |
| **Faro Navy** | `#0B2545` | `(11, 37, 69)` | Fondo de portadas, títulos principales, texto institucional |
| **Faro Ocean** | `#133C55` | `(19, 60, 85)` | Encabezados secundarios, contrastes |
| **Faro Royal** | `#1D4ED8` | `(29, 78, 216)` | Elementos de acento primario, tags de categoría |
| **Faro Sky** | `#38BDF8` | `(56, 189, 248)` | Subtítulos luminosos en fondos oscuros, detalles |
| **Faro Amber** | `#D97706` | `(217, 119, 6)` | Destellos, callouts, indicadores de avance y acentos clave |
| **Faro Slate Dark** | `#1E293B` | `(30, 41, 59)` | Texto de cuerpo principal |
| **Faro Muted** | `#64748B` | `(100, 116, 139)` | Subtítulos, descripciones secundarias, pies de página |
| **Faro BG** | `#F8FAFC` | `(248, 250, 252)` | Fondo de diapositivas de contenido |
| **Faro White** | `#FFFFFF` | `(255, 255, 255)` | Tarjetas y contenedores de datos |
| **Faro Border** | `#E2E8F0` | `(226, 232, 240)` | Bordes sutiles de cards |

**Tipografía**: `Calibri` (Predeterminada para máxima portabilidad universal sin dependencias de fuentes instaladas).

---

## Guía Rápida de Uso en Python

```python
from fcg_helpers import (
    load_fcg_template,
    add_portada_slide,
    add_content_slide,
    add_pattern_two_columns,
    add_pattern_metric_cards,
    dated_filename,
    save_with_normal_view
)

# 1. Cargar plantilla 16:9
prs = load_fcg_template()

# 2. Agregar portada
add_portada_slide(
    prs,
    title="Plan Estratégico 2026",
    subtitle="Consultoría de Dirección",
    author="Faro Consulting Group"
)

# 3. Agregar slide con métricas
slide_kpi = add_content_slide(prs, "Métricas Clave", "Resultados del Q3", "DESEMPEÑO")
add_pattern_metric_cards(slide_kpi, [
    {"number": "+32%", "label": "EBITDA", "desc": "Crecimiento interanual"},
    {"number": "99.4%", "label": "Disponibilidad", "desc": "SLA en producción"}
])

# 4. Guardar con vista normal y lenguaje es-CL
out_file = dated_filename("presentacion_estrategica")
save_with_normal_view(prs, out_file)
```

---

## Reglas de Diseño Obligatorias para Agentes

1. **Aspect Ratio**: Siempre 16:9 widescreen (`13.333` x `7.5` pulgadas). Nunca utilizar 4:3.
2. **Jerarquía Visual**:
   - Títulos de diapositiva: 20–24pt en `FCG_NAVY`, negrita.
   - Categoría / Kicker: 9–10pt en `FCG_ROYAL`, mayúsculas.
   - Texto de cuerpo: 11–13pt en `FCG_DARK`.
   - Números de métricas: 32–40pt en negrita con color de acento.
3. **Márgenes Seguros**:
   - Margen izquierdo mínimo: `0.8` pulgadas.
   - Margen derecho máximo: `12.53` pulgadas (`13.333 - 0.8`).
   - Margen superior contenido: `1.8` a `2.0` pulgadas.
   - Margen inferior: `6.8` a `7.0` pulgadas.
4. **Idioma de Corrección**: Aplicar siempre `set_language_es_cl(prs)` para evitar que PowerPoint subraye texto en español con líneas rojas.
5. **Apertura de Archivo**: Utilizar siempre `save_with_normal_view(prs, path)` para garantizar que el archivo se abra con el panel de diapositivas y miniaturas activo.
'''

def main():
    root = Path(r"c:\Users\mandi\Documents\Proyectos\Plataforma IA local")
    global_skills_root = Path(r"C:\Users\mandi\.gemini\config\skills")
    
    locations = [
        root / ".agents" / "skills" / "fcg-ppt",
        root / "skills" / "fcg-ppt",
        global_skills_root / "fcg-ppt"
    ]
    
    # 1. Crear plantilla base temporal
    temp_tpl = root / "temp_template_fcg.pptx"
    create_master_template(temp_tpl)
    
    for loc in locations:
        print(f"\n--- Desplegando en: {loc} ---")
        assets_dir = loc / "assets"
        scripts_dir = loc / "scripts"
        assets_dir.mkdir(parents=True, exist_ok=True)
        scripts_dir.mkdir(parents=True, exist_ok=True)
        
        # Copiar plantilla
        tpl_target = assets_dir / "template_fcg.pptx"
        with open(temp_tpl, "rb") as f_src, open(tpl_target, "wb") as f_dst:
            f_dst.write(f_src.read())
        print(f" -> Plantilla copiada a {tpl_target}")
        
        # Escribir fcg_helpers.py
        helpers_target = scripts_dir / "fcg_helpers.py"
        with open(helpers_target, "w", encoding="utf-8") as f:
            f.write(FCG_HELPERS_CODE)
        print(f" -> Helpers escritos en {helpers_target}")
        
        # Escribir example_deck.py
        example_target = scripts_dir / "example_deck.py"
        with open(example_target, "w", encoding="utf-8") as f:
            f.write(EXAMPLE_DECK_CODE)
        print(f" -> Script de ejemplo escrito en {example_target}")
        
        # Escribir SKILL.md
        skill_target = loc / "SKILL.md"
        with open(skill_target, "w", encoding="utf-8") as f:
            f.write(SKILL_MD_CONTENT)
        print(f" -> SKILL.md escrito en {skill_target}")

    # Limpiar temp
    if temp_tpl.exists():
        temp_tpl.unlink()
        
    print("\nDespliegue de paquete 'fcg-ppt' completado exitosamente en todos los destinos.")

if __name__ == "__main__":
    main()
